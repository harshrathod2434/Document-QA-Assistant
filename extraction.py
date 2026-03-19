"""
Extraction module: extracts text and images from PDF, DOCX, PPTX, and image files.
Returns structured data with text content, images, and metadata.
"""

import io
import os
from typing import Dict, List, Any

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from PIL import Image


def extract_from_pdf(file_path: str) -> List[Dict[str, Any]]:
    """Extract text and images from a PDF file, page by page."""
    results = []
    filename = os.path.basename(file_path)
    
    try:
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Extract text
            text = page.get_text("text").strip()
            
            # Extract images
            images = []
            image_list = page.get_images(full=True)
            for img_index, img_info in enumerate(image_list):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    if base_image and base_image["image"]:
                        images.append({
                            "data": base_image["image"],
                            "ext": base_image.get("ext", "png"),
                            "label": f"Image {img_index + 1} from page {page_num + 1}"
                        })
                except Exception:
                    continue
            
            results.append({
                "text": text,
                "images": images,
                "metadata": {
                    "source": filename,
                    "file_type": "pdf",
                    "page": page_num + 1,
                    "total_pages": len(doc),
                }
            })
        doc.close()
    except Exception as e:
        results.append({
            "text": f"Error processing PDF {filename}: {str(e)}",
            "images": [],
            "metadata": {"source": filename, "file_type": "pdf", "page": 0}
        })
    
    return results


def extract_from_docx(file_path: str) -> List[Dict[str, Any]]:
    """Extract text and images from a DOCX file."""
    results = []
    filename = os.path.basename(file_path)
    
    try:
        doc = DocxDocument(file_path)
        
        # Extract all paragraph text
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text.strip())
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    full_text.append(" | ".join(row_text))
        
        # Extract images from relationships
        images = []
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_data = rel.target_part.blob
                    ext = rel.target_ref.split(".")[-1].lower() if "." in rel.target_ref else "png"
                    images.append({
                        "data": image_data,
                        "ext": ext,
                        "label": f"Embedded image from {filename}"
                    })
                except Exception:
                    continue
        
        results.append({
            "text": "\n\n".join(full_text),
            "images": images,
            "metadata": {
                "source": filename,
                "file_type": "docx",
                "page": 1,
            }
        })
    except Exception as e:
        results.append({
            "text": f"Error processing DOCX {filename}: {str(e)}",
            "images": [],
            "metadata": {"source": filename, "file_type": "docx", "page": 0}
        })
    
    return results


def extract_from_pptx(file_path: str) -> List[Dict[str, Any]]:
    """Extract text and images from a PPTX file, slide by slide."""
    results = []
    filename = os.path.basename(file_path)
    
    try:
        prs = PptxPresentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            images = []
            
            for shape in slide.shapes:
                # Extract text from shapes
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
                
                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            texts.append(" | ".join(row_text))
                
                # Extract images
                if shape.shape_type == 13:  # Picture type
                    try:
                        image = shape.image
                        image_data = image.blob
                        ext = image.content_type.split("/")[-1] if image.content_type else "png"
                        if ext == "jpeg":
                            ext = "jpg"
                        images.append({
                            "data": image_data,
                            "ext": ext,
                            "label": f"Image from slide {slide_num}"
                        })
                    except Exception:
                        continue
            
            results.append({
                "text": "\n".join(texts),
                "images": images,
                "metadata": {
                    "source": filename,
                    "file_type": "pptx",
                    "slide": slide_num,
                    "total_slides": len(prs.slides),
                }
            })
    except Exception as e:
        results.append({
            "text": f"Error processing PPTX {filename}: {str(e)}",
            "images": [],
            "metadata": {"source": filename, "file_type": "pptx", "slide": 0}
        })
    
    return results


def extract_from_image(file_path: str) -> List[Dict[str, Any]]:
    """Handle standalone image files — pass through as image data for Gemini Vision."""
    filename = os.path.basename(file_path)
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "png"
    
    try:
        with open(file_path, "rb") as f:
            image_data = f.read()
        
        return [{
            "text": "",
            "images": [{
                "data": image_data,
                "ext": ext,
                "label": f"Uploaded image: {filename}"
            }],
            "metadata": {
                "source": filename,
                "file_type": "image",
                "page": 1,
            }
        }]
    except Exception as e:
        return [{
            "text": f"Error reading image {filename}: {str(e)}",
            "images": [],
            "metadata": {"source": filename, "file_type": "image", "page": 0}
        }]


def extract_content(file_path: str) -> List[Dict[str, Any]]:
    """
    Main extraction dispatcher. Detects file type and calls the appropriate extractor.
    Returns a list of page/slide-level results with text, images, and metadata.
    """
    ext = file_path.rsplit(".", 1)[-1].lower() if "." in file_path else ""
    
    extractors = {
        "pdf": extract_from_pdf,
        "docx": extract_from_docx,
        "pptx": extract_from_pptx,
        "png": extract_from_image,
        "jpg": extract_from_image,
        "jpeg": extract_from_image,
    }
    
    extractor = extractors.get(ext)
    if extractor:
        return extractor(file_path)
    else:
        return [{
            "text": f"Unsupported file type: {ext}",
            "images": [],
            "metadata": {"source": os.path.basename(file_path), "file_type": ext, "page": 0}
        }]
